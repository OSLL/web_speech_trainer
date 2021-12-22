from pptx import Presentation


def parse_ppt(presentation_filepath):
    presentation = Presentation(presentation_filepath)
    slides = []
    for slide in presentation.slides:
        slide_info = {
            'title': slide.shapes.title.text if slide.shapes.title else '',
            'words': ''
        }
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_info['words'] += "\n" + shape.text
        slides.append(slide_info)
    return slides